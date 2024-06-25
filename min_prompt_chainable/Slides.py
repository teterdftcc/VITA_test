from pptx import Presentation
from pptx.util import Inches, Pt

def create_slide(prs, layout, title, content):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    return slide

# Create a new presentation
prs = Presentation()

# Set slide layouts
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Mastering the Art of Shrimp Scampi"
subtitle.text = "A Delicious 15-Minute Dinner"

# Slide 2: Hook
create_slide(prs, content_slide_layout, "Why Shrimp Scampi?", 
             "Unlock the secrets to restaurant-quality shrimp scampi in less time than it takes to order takeout!")

# Slide 3: Ingredients
ingredients = """
• 1 pound large shrimp, peeled and deveined
• 4 tablespoons butter
• 4 cloves garlic, minced
• 1/4 cup white wine
• 2 tablespoons lemon juice
• 1/4 cup chopped fresh parsley
• Salt and pepper to taste
• Optional: Red pepper flakes
"""
create_slide(prs, content_slide_layout, "Ingredients", ingredients)

# Slide 4: Cooking Instructions
instructions = """
1. Melt butter in a large skillet over medium heat
2. Add garlic and sauté for 1 minute
3. Add shrimp and cook for 2-3 minutes per side
4. Pour in white wine and lemon juice, simmer for 2 minutes
5. Season with salt, pepper, and optional red pepper flakes
6. Garnish with fresh parsley
"""
create_slide(prs, content_slide_layout, "Cooking Instructions", instructions)

# Slide 5: Tips for Perfect Shrimp Scampi
tips = """
• Use fresh, high-quality shrimp
• Don't overcook the shrimp
• Adjust garlic to your taste preference
• Serve immediately for best flavor
• Pair with crusty bread or pasta
"""
create_slide(prs, content_slide_layout, "Tips for Perfect Shrimp Scampi", tips)

# Slide 6: Variations
variations = """
• Lemon Garlic Shrimp Scampi
• Spicy Shrimp Scampi
• Creamy Shrimp Scampi
• Tomato Basil Shrimp Scampi
• Shrimp Scampi Pasta
"""
create_slide(prs, content_slide_layout, "Variations to Try", variations)

# Slide 7: Conclusion
create_slide(prs, content_slide_layout, "Enjoy Your Homemade Shrimp Scampi!", 
             "With this quick and easy recipe, you can enjoy restaurant-quality shrimp scampi any night of the week!")

# Save the presentation
prs.save('Shrimp_Scampi_Presentation.pptx')
print("Presentation created successfully!")


'''
this is a doc string
This program creates a PowerPoint presentation with seven slides:

1. Title Slide
2. Hook
3. Ingredients
4. Cooking Instructions
5. Tips for Perfect Shrimp Scampi
6. Variations to Try
7. Conclusion

To run this program, you'll need to install the `python-pptx` library:
pip install python-pptx
'''
